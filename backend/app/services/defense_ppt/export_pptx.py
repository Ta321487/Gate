"""python-pptx 导出（永不进学生 ZIP）。"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from app.models import Project

from .check import run_check
from .deck_io import load_deck
from .screenshots import resolve_shot_file
from .themes import ensure_ppt_dirs, ppt_root

_THEME_ACCENT = {
    "scholar": "0B6E75",
    "ink": "1E3A5F",
    "grove": "2F6B4F",
}


def export_pptx(project: Project) -> Path:
    result = run_check(project)
    if not result.get("can_export"):
        errs = [i["message"] for i in result.get("items") or [] if i.get("level") == "error"]
        raise PermissionError(errs[0] if errs else "检查未通过，禁止导出")

    deck = load_deck(project)
    if not deck:
        raise FileNotFoundError("尚无答辩 PPT")

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    theme = str(deck.get("theme") or "scholar")
    hex_accent = _THEME_ACCENT.get(theme, "0B6E75")
    accent = RGBColor(int(hex_accent[0:2], 16), int(hex_accent[2:4], 16), int(hex_accent[4:6], 16))
    ink = RGBColor(0x1A, 0x2B, 0x34)

    def _add_title(slide, text: str, *, top: float = 0.35) -> None:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12), Inches(0.7))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = accent
        p.font.name = "微软雅黑"

    def _add_bullets(slide, bullets: list[dict[str, Any]], *, top: float = 1.3) -> None:
        box = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.5), Inches(5.5))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for b in bullets:
            text = str(b.get("text") if isinstance(b, dict) else b)
            if not text:
                continue
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            para.text = "• " + text
            para.font.size = Pt(18)
            para.font.color.rgb = ink
            para.font.name = "微软雅黑"
            para.space_after = Pt(10)

    def _band(slide) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(0),
            Emu(0),
            Inches(0.35),
            prs.slide_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent
        shape.line.fill.background()

    def _college_master(slide, *, page_title: str = "") -> None:
        """学院母版壳：顶栏 + 页脚（不改业务内容）。"""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(0),
            Emu(0),
            prs.slide_width,
            Inches(0.55),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        head = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12), Inches(0.4))
        hp = head.text_frame.paragraphs[0]
        school = ""
        cover = deck.get("cover") if isinstance(deck.get("cover"), dict) else {}
        school = str(cover.get("school") or "毕业设计答辩")
        hp.text = f"{school} · 答辩文稿"
        hp.font.size = Pt(14)
        hp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        hp.font.name = "微软雅黑"
        foot = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.3))
        fp = foot.text_frame.paragraphs[0]
        fp.text = page_title or str(deck.get("title") or "")
        fp.font.size = Pt(10)
        fp.font.color.rgb = ink
        fp.font.name = "微软雅黑"

    use_master = str(deck.get("master") or "none") == "college_demo"
    pages = [p for p in (deck.get("pages") or []) if isinstance(p, dict)]
    for page in pages:
        slide = prs.slides.add_slide(blank)
        role = page.get("role")
        title = str(page.get("title") or "")
        if use_master:
            _college_master(slide, page_title=title)
        elif str(deck.get("layout_family") or "band") == "band":
            _band(slide)

        title_top = 0.7 if use_master else 0.35
        bullet_top = 1.55 if use_master else 1.3

        if role == "cover":
            cover = page.get("cover") if isinstance(page.get("cover"), dict) else deck.get("cover") or {}
            box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.5), Inches(1.2))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = str(deck.get("title") or project.title or "毕业设计答辩")
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = accent
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
            meta = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10), Inches(3))
            mtf = meta.text_frame
            lines = [
                f"{cover.get('school', '')} · {cover.get('college', '')}",
                f"{cover.get('class_name', '')}  {cover.get('student_name', '')}（{cover.get('student_id', '')}）",
                f"指导教师：{cover.get('advisor', '')}",
            ]
            for i, line in enumerate(lines):
                para = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
                para.text = line
                para.font.size = Pt(16)
                para.font.name = "微软雅黑"
                para.alignment = PP_ALIGN.CENTER
                para.font.color.rgb = ink
            continue

        _add_title(slide, title, top=title_top)

        if role == "toc":
            items = page.get("toc_items") or []
            _add_bullets(
                slide,
                [{"text": f"{i + 1}. {t}"} for i, t in enumerate(items)],
                top=bullet_top,
            )
        elif role in ("modules", "er", "demo"):
            bullets = [b for b in (page.get("bullets") or []) if isinstance(b, dict)]
            _add_bullets(slide, bullets, top=bullet_top)
            fig = page.get("figure") if isinstance(page.get("figure"), dict) else None
            shot = resolve_shot_file(project, fig)
            embedded = False
            if shot and shot.is_file():
                if shot.suffix.lower() in (".png", ".jpg", ".jpeg", ".webp"):
                    try:
                        slide.shapes.add_picture(
                            str(shot), Inches(7.2), Inches(1.5), width=Inches(5.2)
                        )
                        embedded = True
                    except Exception:  # noqa: BLE001
                        embedded = False
            if not embedded and fig:
                note = slide.shapes.add_textbox(Inches(0.9), Inches(5.8), Inches(11), Inches(0.6))
                np = note.text_frame.paragraphs[0]
                label = str(fig.get("label") or "图示嵌自产物")
                if fig.get("path"):
                    label = f"{label} · {fig.get('path')}"
                np.text = label
                np.font.size = Pt(12)
                np.font.name = "微软雅黑"
                np.font.color.rgb = ink
        elif role == "table" or page.get("table"):
            table = page.get("table") if isinstance(page.get("table"), dict) else {}
            headers = list(table.get("headers") or [])
            rows = list(table.get("rows") or [])
            ncols = max(len(headers), max((len(r) for r in rows if isinstance(r, (list, tuple))), default=1))
            nrows = 1 + len(rows)
            if ncols and nrows:
                shape = slide.shapes.add_table(
                    nrows,
                    ncols,
                    Inches(0.9),
                    Inches(1.6 if use_master else 1.4),
                    Inches(11.5),
                    Inches(0.45 * nrows),
                )
                tbl = shape.table
                for c, h in enumerate(headers[:ncols]):
                    cell = tbl.cell(0, c)
                    cell.text = str(h)
                    for para in cell.text_frame.paragraphs:
                        para.font.bold = True
                        para.font.size = Pt(12)
                        para.font.name = "微软雅黑"
                for ri, row in enumerate(rows):
                    vals = list(row) if isinstance(row, (list, tuple)) else [row]
                    for c in range(ncols):
                        cell = tbl.cell(ri + 1, c)
                        cell.text = str(vals[c]) if c < len(vals) else ""
                        for para in cell.text_frame.paragraphs:
                            para.font.size = Pt(11)
                            para.font.name = "微软雅黑"
            bullets = [b for b in (page.get("bullets") or []) if isinstance(b, dict)]
            if bullets:
                _add_bullets(
                    slide,
                    bullets,
                    top=(1.6 if use_master else 1.4) + 0.45 * nrows + 0.2,
                )
        else:
            bullets = [b for b in (page.get("bullets") or []) if isinstance(b, dict)]
            # 溢出：超长要点截断到版心可容纳长度
            clipped = []
            for b in bullets:
                t = str(b.get("text") or "")
                if len(t) > 120:
                    clipped.append({**b, "text": t[:117] + "…"})
                else:
                    clipped.append(b)
            _add_bullets(slide, clipped, top=bullet_top)

    root = ensure_ppt_dirs(project)
    out = root / "export" / "defense.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def export_filename(project: Project) -> str:
    title = (project.title or "defense").strip() or "defense"
    safe = "".join(ch if ch.isalnum() or ch in "-_." else "_" for ch in title)[:60]
    return f"{safe}-答辩.pptx"
